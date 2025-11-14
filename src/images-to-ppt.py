#!/usr/bin/env python3
"""
Images to PowerPoint Converter
Converts a directory of images to a PowerPoint presentation
Each image becomes one slide
"""

import os
import sys
import glob
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Cm
from PIL import Image

# Configuration
CONFIG = {
    'images_dir': os.path.join(os.path.dirname(__file__), '../output/images'),
    'output_dir': os.path.join(os.path.dirname(__file__), '../output'),
    'slide_width': Cm(29.7),  # A4 landscape width
    'slide_height': Cm(21.0), # A4 landscape height
}


def get_image_files(directory):
    """
    Get all PNG image files from directory, sorted by name

    Args:
        directory (str): Directory containing images

    Returns:
        list: Sorted list of image file paths
    """
    pattern = os.path.join(directory, '*.png')
    files = glob.glob(pattern)
    return sorted(files)


def create_ppt_from_images(image_files, output_path, title='HTML Report'):
    """
    Create a PowerPoint presentation from image files

    Args:
        image_files (list): List of image file paths
        output_path (str): Path for output PPTX file
        title (str): Presentation title

    Returns:
        int: Number of slides created
    """
    if not image_files:
        print("❌ No image files found!")
        return 0

    print(f"\n📊 Creating PowerPoint presentation...")
    print(f"  Input: {len(image_files)} images")
    print(f"  Output: {output_path}")

    # Create presentation with custom slide size (A4 landscape)
    prs = Presentation()
    prs.slide_width = CONFIG['slide_width']
    prs.slide_height = CONFIG['slide_height']

    for i, image_path in enumerate(image_files, 1):
        # Use blank slide layout
        blank_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(blank_layout)

        # Get image dimensions
        with Image.open(image_path) as img:
            img_width, img_height = img.size

        # Calculate image aspect ratio
        img_ratio = img_width / img_height
        slide_ratio = float(prs.slide_width) / float(prs.slide_height)

        # Fit image to slide maintaining aspect ratio
        if img_ratio > slide_ratio:
            # Image is wider - fit to width
            width = prs.slide_width
            height = int(width / img_ratio)
            left = 0
            top = int((prs.slide_height - height) / 2)
        else:
            # Image is taller - fit to height
            height = prs.slide_height
            width = int(height * img_ratio)
            left = int((prs.slide_width - width) / 2)
            top = 0

        # Add image to slide
        slide.shapes.add_picture(
            image_path,
            left,
            top,
            width=width,
            height=height
        )

        print(f"  ✓ Slide {i}/{len(image_files)}: {os.path.basename(image_path)}")

    # Save presentation
    prs.save(output_path)
    print(f"\n✅ PowerPoint created successfully!")
    print(f"  Slides: {len(image_files)}")
    print(f"  Size: {os.path.getsize(output_path) / 1024 / 1024:.2f} MB")

    return len(image_files)


def create_presentations_by_report():
    """
    Create separate PPT files for each HTML report
    Groups images by filename prefix
    """
    images_dir = CONFIG['images_dir']
    output_dir = CONFIG['output_dir']

    # Ensure output directory exists
    os.makedirs(output_dir, exist_ok=True)

    # Get all image files
    all_images = get_image_files(images_dir)

    if not all_images:
        print(f"❌ No images found in {images_dir}")
        print("   Please run the HTML to images converter first!")
        return

    print(f"📁 Found {len(all_images)} total images")

    # Group images by report (filename prefix before _page_)
    reports = {}
    for img_path in all_images:
        filename = os.path.basename(img_path)
        # Extract report name (everything before _page_XX.png)
        if '_page_' in filename:
            report_name = filename.split('_page_')[0]
            if report_name not in reports:
                reports[report_name] = []
            reports[report_name].append(img_path)

    print(f"📋 Found {len(reports)} reports:")
    for report_name in reports:
        print(f"  • {report_name}: {len(reports[report_name])} pages")

    # Create PPT for each report
    for report_name, images in reports.items():
        output_path = os.path.join(output_dir, f"{report_name}.pptx")
        create_ppt_from_images(images, output_path, title=report_name)


def create_combined_presentation():
    """
    Create a single PPT file with all images from all reports
    """
    images_dir = CONFIG['images_dir']
    output_dir = CONFIG['output_dir']

    # Ensure output directory exists
    os.makedirs(output_dir, exist_ok=True)

    # Get all image files
    all_images = get_image_files(images_dir)

    if not all_images:
        print(f"❌ No images found in {images_dir}")
        print("   Please run the HTML to images converter first!")
        return

    output_path = os.path.join(output_dir, "combined_reports.pptx")
    create_ppt_from_images(all_images, output_path, title="Combined Reports")


def main():
    """Main execution"""
    print("🚀 Images to PowerPoint Converter Started\n")

    # Check if images directory exists
    if not os.path.exists(CONFIG['images_dir']):
        print(f"❌ Images directory not found: {CONFIG['images_dir']}")
        print("   Please run the HTML to images converter first!")
        sys.exit(1)

    # Create separate PPTs for each report
    print("\n" + "="*60)
    print("Creating separate presentations for each report")
    print("="*60)
    create_presentations_by_report()

    # Also create a combined PPT
    print("\n" + "="*60)
    print("Creating combined presentation")
    print("="*60)
    create_combined_presentation()

    print("\n✨ All presentations created successfully!")
    print(f"📁 Output directory: {CONFIG['output_dir']}")


if __name__ == '__main__':
    main()
